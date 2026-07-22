from pathlib import Path

import pandas as pd
import pdfplumber
import pytesseract
from docx import Document
from pdf2image import convert_from_path
from PIL import Image
from pptx import Presentation

from .files import clean_text


class ExtractionError(Exception):
    """Raised when text extraction cannot produce usable content."""


def extract_pdf_text(path: str) -> tuple[str, bool]:
    text_parts = []
    scanned_pages = 0
    try:
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text() or ""
                if len(page_text.strip()) < 20:
                    scanned_pages += 1
                text_parts.append(page_text)
            normal_text = "\n".join(text_parts).strip()
            is_scanned = bool(pdf.pages) and scanned_pages >= max(1, len(pdf.pages) // 2)
    except Exception as exc:
        raise ExtractionError(f"Unable to read PDF: {exc}") from exc

    if normal_text and not is_scanned:
        return clean_text(normal_text), False

    try:
        images = convert_from_path(path, dpi=300)
        ocr_text = "\n".join(pytesseract.image_to_string(image) for image in images)
        return clean_text(ocr_text), True
    except Exception as exc:
        if normal_text:
            return clean_text(normal_text), True
        raise ExtractionError(
            "This PDF appears to be scanned, but OCR could not run. Install Poppler and Tesseract, then try again."
        ) from exc


def extract_docx_text(path: str) -> str:
    try:
        document = Document(path)
        paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        table_text = []
        for table in document.tables:
            for row in table.rows:
                table_text.append(" | ".join(cell.text.strip() for cell in row.cells if cell.text.strip()))
        return clean_text("\n".join(paragraphs + table_text))
    except Exception as exc:
        raise ExtractionError(f"Unable to read Word document: {exc}") from exc


def extract_pptx_text(path: str) -> str:
    try:
        presentation = Presentation(path)
        chunks = []
        for index, slide in enumerate(presentation.slides, start=1):
            slide_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text.strip())
            if slide_parts:
                chunks.append(f"Slide {index}: " + "\n".join(slide_parts))
        return clean_text("\n".join(chunks))
    except Exception as exc:
        raise ExtractionError(f"Unable to read PowerPoint file: {exc}") from exc


def extract_csv_text(path: str) -> str:
    try:
        dataframe = pd.read_csv(path, encoding_errors="ignore")
        preview = dataframe.to_string(index=False, max_rows=500)
        summary = f"Columns: {', '.join(map(str, dataframe.columns))}\nRows: {len(dataframe)}\n\n{preview}"
        return clean_text(summary)
    except Exception as exc:
        raise ExtractionError(f"Unable to read CSV file: {exc}") from exc


def extract_txt_text(path: str) -> str:
    encodings = ["utf-8", "utf-16", "latin-1"]
    for encoding in encodings:
        try:
            with open(path, "r", encoding=encoding) as handle:
                return clean_text(handle.read())
        except UnicodeDecodeError:
            continue
        except Exception as exc:
            raise ExtractionError(f"Unable to read text file: {exc}") from exc
    raise ExtractionError("Unable to decode text file.")


def extract_image_text(path: str) -> str:
    try:
        with Image.open(path) as image:
            return clean_text(pytesseract.image_to_string(image))
    except Exception as exc:
        raise ExtractionError("Unable to run OCR on image. Install Tesseract and try again.") from exc


def extract_text(path: str, filename: str | None = None) -> tuple[str, dict]:
    source_name = filename or path
    extension = Path(source_name).suffix.lower()
    metadata = {"source": source_name, "extension": extension, "used_ocr": False}

    if extension == ".pdf":
        text, used_ocr = extract_pdf_text(path)
        metadata["used_ocr"] = used_ocr
    elif extension == ".docx":
        text = extract_docx_text(path)
    elif extension == ".pptx":
        text = extract_pptx_text(path)
    elif extension == ".txt":
        text = extract_txt_text(path)
    elif extension == ".csv":
        text = extract_csv_text(path)
    elif extension in {".jpg", ".jpeg", ".png", ".webp"}:
        text = extract_image_text(path)
        metadata["used_ocr"] = True
    else:
        raise ExtractionError("Unsupported file format.")

    if not text or len(text.strip()) < 5:
        raise ExtractionError("No readable text was found in this file.")
    metadata["characters"] = len(text)
    metadata["words"] = len(text.split())
    return text, metadata
